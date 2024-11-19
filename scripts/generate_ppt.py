from pptx import Presentation
from pptx.util import Inches

# Initialize the presentation
prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Fantasy Premier League Optimization"
subtitle.text = "Data Product Design and Optimization Process"

# Slide 1: Objectives
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Objectives"
content.text = (
    "1. Build the best fantasy team each week to maximize points.\n"
    "2. Suggest weekly transfers to improve team performance.\n"
    "3. Combine data from football and fantasy APIs for decision-making.\n"
    "4. Use optimization models to meet FPL constraints and rules."
)

# Slide 2: System Architecture
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "System Architecture"
content.text = (
    "1. Data Sources:\n"
    "   - Football API: Player stats, fixtures, form, injuries.\n"
    "   - Fantasy API: Prices, ownership, points, team structure.\n"
    "2. Data Processing:\n"
    "   - Calculate projected points for each player.\n"
    "3. Optimization Engine:\n"
    "   - Select optimal team or suggest transfers.\n"
    "4. Outputs:\n"
    "   - Best team.\n"
    "   - Suggested weekly changes."
)

# Slide 3: Data Integration
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Data Integration"
content.text = (
    "1. Extract data from Football API:\n"
    "   - Player stats: goals, assists, xG, xA, clean sheets.\n"
    "   - Fixtures: Match difficulty and schedule.\n"
    "2. Extract data from Fantasy API:\n"
    "   - Prices, positions, ownership stats, past points.\n"
    "3. Merge data using unique player IDs or names."
)

# Slide 4: Projected Points Calculation
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Projected Points Calculation"
content.text = (
    "1. Use player stats and upcoming fixtures to calculate expected points.\n"
    "2. Factors to consider:\n"
    "   - Player form: Recent performance metrics.\n"
    "   - Fixture difficulty: Opponent strength.\n"
    "   - Injury/rotation risk: Penalize uncertain players.\n"
    "   - Goal/assist/clean sheet potential.\n"
    "3. Example formula:\n"
    "   Projected Points = Form Weight × Recent Points + Fixture Weight × Difficulty + xG/xA Contributions"
)

# Slide 5: Optimization Model
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Optimization Model"
content.text = (
    "1. Objective:\n"
    "   - Maximize total projected points.\n"
    "2. Constraints:\n"
    "   - Budget: ≤ £100M.\n"
    "   - Squad size: 15 players.\n"
    "   - Positions: 2 GKs, 5 DEFs, 5 MIDs, 3 FWDs.\n"
    "   - Team limit: Max 3 players per team.\n"
    "3. Solve using linear or integer programming."
)

# Slide 6: Outputs
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Outputs"
content.text = (
    "1. Optimal Team:\n"
    "   - List of players with positions, prices, and projected points.\n"
    "2. Transfer Suggestions:\n"
    "   - Players to transfer in and out.\n"
    "   - Impact on projected points.\n"
    "3. Customization:\n"
    "   - Allow user preferences (e.g., avoid certain teams)."
)

# Slide 7: Next Steps
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Next Steps"
content.text = (
    "1. Complete data extraction and integration.\n"
    "2. Develop projected points calculation.\n"
    "3. Build and test the optimization model.\n"
    "4. Create a user interface for actionable outputs.\n"
    "5. Test weekly and refine based on performance."
)

# Save the presentation
output_file = "Fantasy_Premier_League_Optimization.pptx"
prs.save(output_file)
print(f"Presentation saved as {output_file}")
